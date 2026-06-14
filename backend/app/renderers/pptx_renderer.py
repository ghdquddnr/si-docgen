"""화면정의서 PPT 렌더러.

순수 함수: (검증된 ScreenSpecDocument, 템플릿 경로) → 출력 .pptx 파일.
템플릿의 표준 슬라이드(2번째)를 XML 레벨로 복제해 화면 수만큼 만들고,
shape name 기반으로 텍스트/표 값만 주입한다. 서식은 템플릿의 책임이다.
"""

import copy
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.presentation import Presentation as PresentationType
from pptx.shapes.base import BaseShape
from pptx.slide import Slide
from pptx.text.text import TextFrame
from pptx.util import Pt

from app.exceptions import RenderError
from app.schemas.screen_spec import Screen, ScreenField, ScreenSpecDocument

# 템플릿 구조 상수 (backend/templates/screen_spec.pptx 와 일치해야 함)
COVER_SLIDE_INDEX = 0
STANDARD_SLIDE_INDEX = 1
FIELD_TABLE_COLUMNS = 5


def circled_number(no: int) -> str:
    """항목 번호를 ①②③ 형태의 원문자로 변환한다 (1~20 지원)."""
    if 1 <= no <= 20:
        return chr(0x2460 + no - 1)
    return str(no)


def render_screen_spec(doc: ScreenSpecDocument, template_path: Path, output_path: Path) -> Path:
    """표지를 채우고 표준 슬라이드를 화면 수만큼 복제해 화면정의서를 생성한다.

    각 화면의 목업은 항목 정의에서 **편집 가능한 PPT 도형**(번호 배지·라벨·입력/버튼 도형)으로
    그린다 — 이미지가 아니므로 사용자가 PowerPoint 에서 위치·크기·텍스트를 직접 수정할 수 있다.
    """
    if not template_path.is_file():
        raise RenderError(f"템플릿 파일이 없습니다: {template_path}")

    prs = Presentation(str(template_path))
    if len(prs.slides) < 2:
        raise RenderError(f"템플릿에 표지/표준 슬라이드 2장이 필요합니다: {template_path}")

    _fill_cover(prs.slides[COVER_SLIDE_INDEX], doc)

    standard = prs.slides[STANDARD_SLIDE_INDEX]
    for screen in doc.screens:
        new_slide = _duplicate_slide(prs, standard)
        _fill_screen(new_slide, screen)
        _draw_mockup(new_slide, screen)
    _remove_slide(prs, STANDARD_SLIDE_INDEX)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return output_path


def _duplicate_slide(prs: PresentationType, source: Slide) -> Slide:
    """슬라이드의 모든 shape 를 XML deepcopy 해 새 슬라이드로 복제한다."""
    dest = prs.slides.add_slide(source.slide_layout)
    for shape in source.shapes:
        dest.shapes._spTree.insert_element_before(copy.deepcopy(shape._element), "p:extLst")
    return dest


def _remove_slide(prs: PresentationType, index: int) -> None:
    """슬라이드 목록에서 index 위치의 슬라이드를 제거한다."""
    sld_id_lst = prs.slides._sldIdLst
    sld_id_lst.remove(list(sld_id_lst)[index])


def _find_shape(slide: Slide, name: str) -> BaseShape:
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    raise RenderError(f"슬라이드에 shape '{name}' 가 없습니다 (템플릿 구조 확인 필요)")


def _set_text(text_frame: TextFrame, text: str) -> None:
    """첫 run 의 서식을 유지한 채 텍스트만 교체한다 ('\\n' 은 줄바꿈으로 변환)."""
    para = text_frame.paragraphs[0]
    if not para.runs:
        raise RenderError("템플릿 텍스트 프레임에 서식 기준 run 이 없습니다")
    para.runs[0].text = text
    for run in para.runs[1:]:
        run._r.getparent().remove(run._r)


def _fill_cover(slide: Slide, doc: ScreenSpecDocument) -> None:
    values = {
        "cover_project": doc.project_name,
        "cover_system": doc.system_name,
        "cover_author": doc.author,
        "cover_date": doc.written_date.isoformat(),
    }
    for name, value in values.items():
        _set_text(_find_shape(slide, name).text_frame, value)


def _fill_screen(slide: Slide, screen: Screen) -> None:
    _set_text(_find_shape(slide, "slide_title").text_frame, screen.screen_name)
    _set_text(_find_shape(slide, "screen_id").text_frame, screen.screen_id)
    _set_text(_find_shape(slide, "screen_name").text_frame, screen.screen_name)
    _set_text(_find_shape(slide, "menu_path").text_frame, screen.menu_path)

    logic_text = "\n".join(f"{i}. {line}" for i, line in enumerate(screen.logic, start=1))
    _set_text(_find_shape(slide, "logic_text").text_frame, logic_text)

    _fill_field_table(slide, screen)


# 목업 와이어프레임 색 (slate/indigo 계열)
_WIRE_BORDER = RGBColor(0x94, 0xA3, 0xB8)
_INPUT_FILL = RGBColor(0xFF, 0xFF, 0xFF)
_BUTTON_FILL = RGBColor(0xE2, 0xE8, 0xF0)
_BADGE_FILL = RGBColor(0x4F, 0x46, 0xE5)
_LABEL_COLOR = RGBColor(0x33, 0x41, 0x55)
_WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def _ftype_has(field_type: str, *keywords: str) -> bool:
    lowered = field_type.lower()
    return any(k in lowered for k in keywords)


def _style_text(
    text_frame: TextFrame,
    text: str,
    size_pt: int,
    color: RGBColor,
    *,
    align: PP_ALIGN | None = None,
) -> None:
    """도형 텍스트를 작은 폰트로 채운다 (목업 와이어프레임용)."""
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    text_frame.margin_left = text_frame.margin_right = Pt(2)
    text_frame.margin_top = text_frame.margin_bottom = Pt(1)
    para = text_frame.paragraphs[0]
    para.text = text
    if align is not None:
        para.alignment = align
    run = para.runs[0]
    run.font.size = Pt(size_pt)
    run.font.color.rgb = color


def _draw_mockup(slide: Slide, screen: Screen) -> None:
    """mockup_area 안에 항목 목록을 편집 가능한 도형(번호·라벨·컨트롤)으로 그린다."""
    area = _find_shape(slide, "mockup_area")
    _set_text(area.text_frame, "")  # 자리표시 문구 제거 (테두리 사각형은 액자로 유지)
    if not screen.fields:
        return

    pad = int(area.width * 0.05)
    inner_left = area.left + pad
    inner_top = area.top + pad
    inner_w = area.width - 2 * pad
    inner_h = area.height - 2 * pad

    n = len(screen.fields)
    row_h = int(inner_h / n)
    ctrl_h = int(min(row_h * 0.55, Pt(22)))
    badge_d = int(min(row_h * 0.5, Pt(16)))
    gap = int(inner_w * 0.02)
    label_w = int(inner_w * 0.34)
    ctrl_left = inner_left + badge_d + gap + label_w + gap
    ctrl_w = inner_left + inner_w - ctrl_left

    for i, field in enumerate(screen.fields):
        row_top = inner_top + i * row_h
        _draw_field_row(
            slide,
            field,
            row_top,
            row_h,
            inner_left,
            badge_d,
            gap,
            label_w,
            ctrl_left,
            ctrl_w,
            ctrl_h,
            inner_w,
        )


def _draw_field_row(
    slide: Slide,
    field: ScreenField,
    row_top: int,
    row_h: int,
    inner_left: int,
    badge_d: int,
    gap: int,
    label_w: int,
    ctrl_left: int,
    ctrl_w: int,
    ctrl_h: int,
    inner_w: int,
) -> None:
    # ── 번호 배지 (목업의 ①②③ 대신 PowerPoint 도형 + 숫자)
    badge_top = row_top + (row_h - badge_d) // 2
    badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, inner_left, badge_top, badge_d, badge_d)
    badge.fill.solid()
    badge.fill.fore_color.rgb = _BADGE_FILL
    badge.line.fill.background()
    badge.name = f"mockup_badge_{field.no}"
    _style_text(badge.text_frame, str(field.no), 8, _WHITE, align=PP_ALIGN.CENTER)

    # ── 항목 라벨 (필수 항목은 * 표시)
    label = slide.shapes.add_textbox(inner_left + badge_d + gap, row_top, label_w, row_h)
    label.name = f"mockup_label_{field.no}"
    _style_text(label.text_frame, field.name + (" *" if field.required else ""), 9, _LABEL_COLOR)

    # ── 컨트롤 (유형별 도형)
    ctrl_top = row_top + (row_h - ctrl_h) // 2
    if _ftype_has(field.field_type, "버튼", "button"):
        width = min(ctrl_w, int(inner_w * 0.28))
        ctrl = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, ctrl_left, ctrl_top, width, ctrl_h
        )
        ctrl.fill.solid()
        ctrl.fill.fore_color.rgb = _BUTTON_FILL
        ctrl.line.color.rgb = _WIRE_BORDER
        _style_text(ctrl.text_frame, field.name, 8, _LABEL_COLOR, align=PP_ALIGN.CENTER)
    elif _ftype_has(field.field_type, "체크", "check", "라디오", "radio"):
        ctrl = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, ctrl_left, ctrl_top, ctrl_h, ctrl_h)
        ctrl.fill.solid()
        ctrl.fill.fore_color.rgb = _INPUT_FILL
        ctrl.line.color.rgb = _WIRE_BORDER
    else:
        ctrl = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, ctrl_left, ctrl_top, ctrl_w, ctrl_h)
        ctrl.fill.solid()
        ctrl.fill.fore_color.rgb = _INPUT_FILL
        ctrl.line.color.rgb = _WIRE_BORDER
        if _ftype_has(field.field_type, "콤보", "combo", "드롭", "select", "셀렉"):
            _style_text(ctrl.text_frame, "▼", 8, _WIRE_BORDER, align=PP_ALIGN.RIGHT)
    ctrl.name = f"mockup_ctrl_{field.no}"


def _fill_field_table(slide: Slide, screen: Screen) -> None:
    """항목 정의 표의 서식 기준 행(2행)을 복제해 항목 수만큼 채운다."""
    frame = _find_shape(slide, "field_table")
    if not frame.has_table:
        raise RenderError("shape 'field_table' 이 표가 아닙니다 (템플릿 구조 확인 필요)")
    table = frame.table
    if len(table.columns) != FIELD_TABLE_COLUMNS or len(table.rows) < 2:
        raise RenderError("항목 정의 표는 5열 + 헤더/서식 기준 행 구조여야 합니다")

    tbl = table._tbl
    sample_tr = tbl.tr_lst[1]
    for _ in range(len(screen.fields) - 1):
        tbl.append(copy.deepcopy(sample_tr))

    # python-pptx 의 행 컬렉션은 슬라이싱을 지원하지 않아 list 로 변환한다
    for row, field in zip(list(table.rows)[1:], screen.fields, strict=True):
        values = [
            circled_number(field.no),
            field.name,
            field.field_type,
            "Y" if field.required else "N",
            field.description,
        ]
        for cell, value in zip(row.cells, values, strict=True):
            _set_text(cell.text_frame, value)
