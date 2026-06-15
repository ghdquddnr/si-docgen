"""제안서 PPT 렌더러.

순수 함수: (검증된 ProposalDocument, 템플릿 경로) → 출력 .pptx 파일.
템플릿의 표지를 채우고, 표준 내용 슬라이드(2번째)를 슬라이드 수만큼 XML 복제해
shape name 기반으로 제목·불릿만 주입한다. 서식은 템플릿의 책임이다(절대 원칙 4).
불릿은 편집 가능한 텍스트 프레임에 채워 사용자가 PowerPoint 에서 직접 수정할 수 있다.
"""

import copy
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.presentation import Presentation as PresentationType
from pptx.shapes.base import BaseShape
from pptx.slide import Slide
from pptx.text.text import TextFrame
from pptx.util import Pt

from app.exceptions import RenderError
from app.schemas.proposal import ProposalDocument, ProposalSlide

COVER_SLIDE_INDEX = 0
STANDARD_SLIDE_INDEX = 1

_FONT_NAME = "맑은 고딕"
_BULLET_COLOR = RGBColor(0x33, 0x41, 0x55)
_BULLET_SIZE = Pt(18)


def render_proposal(doc: ProposalDocument, template_path: Path, output_path: Path) -> Path:
    """표지를 채우고 표준 슬라이드를 제안 슬라이드 수만큼 복제해 제안서를 생성한다."""
    if not template_path.is_file():
        raise RenderError(f"템플릿 파일이 없습니다: {template_path}")

    prs = Presentation(str(template_path))
    if len(prs.slides) < 2:
        raise RenderError(f"템플릿에 표지/표준 슬라이드 2장이 필요합니다: {template_path}")

    _fill_cover(prs.slides[COVER_SLIDE_INDEX], doc)

    standard = prs.slides[STANDARD_SLIDE_INDEX]

    # 목차 슬라이드 — 섹션 제목에서 자동 구성(렌더러 책임, 항상 본문과 일치)
    toc = _duplicate_slide(prs, standard)
    _set_text(_find_shape(toc, "slide_title").text_frame, "목차")
    toc_lines = [f"{i}. {s.title}" for i, s in enumerate(doc.slides, start=1)]
    _set_bullets(_find_shape(toc, "slide_body").text_frame, toc_lines, marker="")

    for slide_data in doc.slides:
        new_slide = _duplicate_slide(prs, standard)
        _fill_slide(new_slide, slide_data)
    _remove_slide(prs, STANDARD_SLIDE_INDEX)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return output_path


def _duplicate_slide(prs: PresentationType, source: Slide) -> Slide:
    dest = prs.slides.add_slide(source.slide_layout)
    for shape in source.shapes:
        dest.shapes._spTree.insert_element_before(copy.deepcopy(shape._element), "p:extLst")
    return dest


def _remove_slide(prs: PresentationType, index: int) -> None:
    sld_id_lst = prs.slides._sldIdLst
    sld_id_lst.remove(list(sld_id_lst)[index])


def _find_shape(slide: Slide, name: str) -> BaseShape:
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    raise RenderError(f"슬라이드에 shape '{name}' 가 없습니다 (템플릿 구조 확인 필요)")


def _set_text(text_frame: TextFrame, text: str) -> None:
    """첫 run 의 서식을 유지한 채 텍스트만 교체한다."""
    para = text_frame.paragraphs[0]
    if not para.runs:
        raise RenderError("템플릿 텍스트 프레임에 서식 기준 run 이 없습니다")
    para.runs[0].text = text
    for run in para.runs[1:]:
        run._r.getparent().remove(run._r)


def _korean_font(run, size: Pt, color: RGBColor) -> None:  # noqa: ANN001
    run.font.name = _FONT_NAME
    run.font.size = size
    run.font.color.rgb = color
    rpr = run._r.get_or_add_rPr()
    ea = rpr.find(qn("a:ea"))
    if ea is None:
        ea = etree.SubElement(rpr, qn("a:ea"))
    ea.set("typeface", _FONT_NAME)


def _set_bullets(text_frame: TextFrame, lines: list[str], marker: str = "• ") -> None:
    """본문 텍스트 프레임을 문단으로 채운다 (각 줄 = 한 문단, marker 머리표).

    marker 기본은 불릿('• '). 목차처럼 번호를 직접 넣을 때는 marker="" 로 호출한다.
    """
    text_frame.clear()  # 문단 1개만 남는다
    for i, line in enumerate(lines):
        para = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
        run = para.add_run()
        run.text = f"{marker}{line}"
        _korean_font(run, _BULLET_SIZE, _BULLET_COLOR)


def _fill_cover(slide: Slide, doc: ProposalDocument) -> None:
    values = {
        "cover_title": doc.title,
        "cover_project": doc.project_name,
        "cover_client": doc.client,
        "cover_proposer": doc.author,
        "cover_date": doc.written_date.isoformat(),
    }
    for name, value in values.items():
        _set_text(_find_shape(slide, name).text_frame, value)


def _fill_slide(slide: Slide, data: ProposalSlide) -> None:
    _set_text(_find_shape(slide, "slide_title").text_frame, data.title)
    _set_bullets(_find_shape(slide, "slide_body").text_frame, data.bullets)
