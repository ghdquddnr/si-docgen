"""양식 보관함 — 사용자 업로드 템플릿의 저장·검증·조회·경로 해석.

설계(B1, 구조 보존형): 사용자는 기본 양식을 내려받아 로고·색·결재란 등 **서식만** 바꾸고
구조(시트·docxtpl 태그·도형 이름)는 유지한 채 재업로드한다. 업로드 시 기본 양식과
구조 마커가 호환되는지 검증해, 호환되지 않으면 거부한다. 렌더러는 순수 함수를 유지하며
(절대 원칙 2), 이 서비스가 종류별로 사용할 템플릿 경로만 골라 전달한다.
"""

import io
import logging
import uuid
from pathlib import Path
from typing import BinaryIO

from docxtpl import DocxTemplate
from openpyxl import load_workbook
from pptx import Presentation
from sqlalchemy import select
from sqlalchemy.orm import Session

from app.config import get_settings
from app.db.models import Template, TemplateFolder
from app.exceptions import SiDocgenError
from app.pipelines.generate_chain import REQUIREMENT_SPEC_TEMPLATE, SCREEN_SPEC_TEMPLATE
from app.pipelines.generate_interface_spec import INTERFACE_SPEC_TEMPLATE
from app.pipelines.generate_proposal import PROPOSAL_TEMPLATE
from app.pipelines.generate_table_spec import TABLE_SPEC_TEMPLATE
from app.pipelines.generate_test_scenario import RTM_TEMPLATE, TEST_SCENARIO_TEMPLATE
from app.pipelines.generate_user_manual import USER_MANUAL_TEMPLATE
from app.pipelines.generate_wbs import WBS_TEMPLATE

logger = logging.getLogger(__name__)

# 산출물 종류 → 기본(번들) 양식 경로. 사용자 양식 검증·미선택 시 폴백의 기준.
DEFAULT_TEMPLATES: dict[str, Path] = {
    "proposal": PROPOSAL_TEMPLATE,
    "test_scenario": TEST_SCENARIO_TEMPLATE,
    "rtm": RTM_TEMPLATE,
    "requirement_spec": REQUIREMENT_SPEC_TEMPLATE,
    "screen_spec": SCREEN_SPEC_TEMPLATE,
    "wbs": WBS_TEMPLATE,
    "table_spec": TABLE_SPEC_TEMPLATE,
    "interface_spec": INTERFACE_SPEC_TEMPLATE,
    "user_manual": USER_MANUAL_TEMPLATE,
}

# 종류 → 파일 확장자
KIND_EXT: dict[str, str] = {
    "proposal": ".pptx",
    "test_scenario": ".xlsx",
    "rtm": ".xlsx",
    "wbs": ".xlsx",
    "table_spec": ".xlsx",
    "interface_spec": ".xlsx",
    "requirement_spec": ".docx",
    "user_manual": ".docx",
    "screen_spec": ".pptx",
}

# 사람이 읽는 종류 이름 (UI 라벨)
KIND_LABELS: dict[str, str] = {
    "proposal": "제안서",
    "test_scenario": "테스트시나리오",
    "rtm": "요건추적표(RTM)",
    "requirement_spec": "요구사항정의서",
    "screen_spec": "화면정의서",
    "wbs": "WBS",
    "table_spec": "테이블정의서",
    "interface_spec": "인터페이스정의서",
    "user_manual": "사용자 매뉴얼",
}


class TemplateValidationError(SiDocgenError):
    """업로드 양식이 기본 양식과 구조 호환되지 않을 때 발생 (API 400)."""


class TemplateNotFoundError(SiDocgenError):
    """존재하지 않는 템플릿/폴더 참조 (API 404)."""


def templates_dir() -> Path:
    return Path(get_settings().templates_dir)


def template_path(template_id: str, kind: str) -> Path:
    return templates_dir() / f"{template_id}{KIND_EXT[kind]}"


# ── 구조 마커 추출 (검증 기준은 기본 양식에서 자동 도출) ──────────────────────


def _xlsx_sheets(src: Path | BinaryIO) -> set[str]:
    wb = load_workbook(src, read_only=True)
    try:
        return set(wb.sheetnames)
    finally:
        wb.close()


def _docx_vars(src: Path | BinaryIO) -> set[str]:
    arg = str(src) if isinstance(src, Path) else src
    return set(DocxTemplate(arg).get_undeclared_template_variables())


def _pptx_shapes(src: Path | BinaryIO) -> set[str]:
    arg = str(src) if isinstance(src, Path) else src
    names: set[str] = set()
    for slide in Presentation(arg).slides:
        for shape in slide.shapes:
            if shape.name:
                names.add(shape.name)
    return names


def validate_template(kind: str, data: bytes, suffix: str) -> None:
    """업로드 양식(바이트)이 해당 종류의 기본 양식과 구조 호환되는지 검증한다.

    호환 기준은 기본 양식에서 자동 도출한다(xlsx=시트명, docx=docxtpl 변수, pptx=도형명).
    누락·형식 불일치가 있으면 TemplateValidationError 를 올린다. (디스크 쓰기 전에 메모리에서 검증)
    """
    if kind not in DEFAULT_TEMPLATES:
        raise TemplateValidationError(f"지원하지 않는 양식 종류: {kind}")
    ext = KIND_EXT[kind]
    if suffix.lower() != ext:
        raise TemplateValidationError(
            f"{KIND_LABELS[kind]} 양식은 {ext} 파일이어야 합니다 (받은 형식: {suffix or '없음'})"
        )
    default = DEFAULT_TEMPLATES[kind]
    try:
        if ext == ".xlsx":
            required, got = _xlsx_sheets(default), _xlsx_sheets(io.BytesIO(data))
            label = "시트"
        elif ext == ".docx":
            required, got = _docx_vars(default), _docx_vars(io.BytesIO(data))
            label = "치환 태그"
        else:
            required, got = _pptx_shapes(default), _pptx_shapes(io.BytesIO(data))
            label = "도형"
    except TemplateValidationError:
        raise
    except Exception as exc:  # 파일이 손상됐거나 형식이 맞지 않음
        raise TemplateValidationError(
            f"양식 파일을 열 수 없습니다 ({ext} 형식인지 확인하세요): {exc}"
        ) from exc

    missing = required - got
    if missing:
        raise TemplateValidationError(
            f"{KIND_LABELS[kind]} 기본 양식의 필수 {label} 가 빠져 있습니다: "
            f"{', '.join(sorted(missing))}. 기본 양식을 받아 서식만 수정해 다시 올려주세요."
        )


# ── 저장 / 조회 / 삭제 ────────────────────────────────────────────────────


def save_template(
    db: Session, *, name: str, kind: str, folder_id: str | None, filename: str, data: bytes
) -> Template:
    """업로드 양식을 검증 후 저장한다. 폴더가 지정되면 존재해야 한다."""
    if kind not in DEFAULT_TEMPLATES:
        raise TemplateValidationError(f"지원하지 않는 양식 종류: {kind}")
    if folder_id is not None and db.get(TemplateFolder, folder_id) is None:
        raise TemplateNotFoundError(f"폴더를 찾을 수 없습니다: {folder_id}")

    validate_template(kind, data, Path(filename).suffix)  # 디스크 쓰기 전에 검증

    template_id = uuid.uuid4().hex
    target = template_path(template_id, kind)
    target.parent.mkdir(parents=True, exist_ok=True)
    target.write_bytes(data)

    tpl = Template(
        id=template_id,
        name=name or filename,
        kind=kind,
        folder_id=folder_id,
        original_filename=filename,
    )
    db.add(tpl)
    db.commit()
    db.refresh(tpl)
    logger.info("양식 저장: id=%s kind=%s name=%s", template_id, kind, tpl.name)
    return tpl


def delete_template(db: Session, template_id: str) -> bool:
    tpl = db.get(Template, template_id)
    if tpl is None:
        return False
    template_path(template_id, tpl.kind).unlink(missing_ok=True)
    db.delete(tpl)
    db.commit()
    return True


def create_folder(db: Session, *, name: str, parent_id: str | None) -> TemplateFolder:
    if parent_id is not None and db.get(TemplateFolder, parent_id) is None:
        raise TemplateNotFoundError(f"상위 폴더를 찾을 수 없습니다: {parent_id}")
    folder = TemplateFolder(id=uuid.uuid4().hex, name=name, parent_id=parent_id)
    db.add(folder)
    db.commit()
    db.refresh(folder)
    return folder


def delete_folder(db: Session, folder_id: str) -> bool:
    """폴더와 그 하위(폴더·양식)를 재귀 삭제한다."""
    folder = db.get(TemplateFolder, folder_id)
    if folder is None:
        return False
    for child in db.scalars(
        select(TemplateFolder).where(TemplateFolder.parent_id == folder_id)
    ).all():
        delete_folder(db, child.id)
    for tpl in db.scalars(select(Template).where(Template.folder_id == folder_id)).all():
        delete_template(db, tpl.id)
    db.delete(folder)
    db.commit()
    return True


def list_folders(db: Session) -> list[TemplateFolder]:
    return list(db.scalars(select(TemplateFolder).order_by(TemplateFolder.name)).all())


def list_templates(db: Session, kind: str | None = None) -> list[Template]:
    stmt = select(Template).order_by(Template.created_at.desc())
    if kind is not None:
        stmt = stmt.where(Template.kind == kind)
    return list(db.scalars(stmt).all())


def resolve_template_path(db: Session, kind: str, template_id: str | None) -> Path:
    """종류·선택 id 로 실제 사용할 양식 경로를 해석한다. 미선택/없으면 기본 양식."""
    if template_id:
        tpl = db.get(Template, template_id)
        if tpl is not None and tpl.kind == kind:
            path = template_path(template_id, kind)
            if path.is_file():
                return path
        logger.warning(
            "선택한 양식을 찾을 수 없어 기본 양식 사용: kind=%s id=%s", kind, template_id
        )
    return DEFAULT_TEMPLATES[kind]


def resolve_all(db: Session, template_ids: dict | None) -> dict[str, Path]:
    """잡의 template_ids 맵을 종류→경로 맵으로 해석한다(미지정 종류는 기본 양식)."""
    ids = template_ids or {}
    return {kind: resolve_template_path(db, kind, ids.get(kind)) for kind in DEFAULT_TEMPLATES}
