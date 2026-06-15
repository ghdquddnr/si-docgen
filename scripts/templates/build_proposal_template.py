"""제안서 PPT 양식(backend/templates/proposal.pptx) 제작 스크립트.

표지 슬라이드 + 내용 표준 슬라이드 1장을 구성한다. 렌더러는 표준 슬라이드를 슬라이드 수만큼
복제해 제목·불릿만 주입하며, 서식은 이 템플릿이 보존한다(절대 원칙 4). 일회성 제작 도구.

실행: uv run python scripts/templates/build_proposal_template.py

shape name 규약:
- 표지: cover_title / cover_project / cover_client / cover_proposer / cover_date
- 표준 슬라이드: slide_title / slide_body
"""

from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.slide import Slide
from pptx.util import Inches, Pt

OUTPUT_PATH = Path(__file__).resolve().parents[2] / "backend" / "templates" / "proposal.pptx"

FONT_NAME = "맑은 고딕"
INDIGO = RGBColor(0x4F, 0x46, 0xE5)
SLATE = RGBColor(0x33, 0x41, 0x55)
GRAY_TEXT = RGBColor(0x80, 0x80, 0x80)
BLACK = RGBColor(0x1E, 0x29, 0x3B)


def set_run_font(run, size: float, *, bold: bool = False, color: RGBColor = BLACK) -> None:  # noqa: ANN001
    run.font.name = FONT_NAME
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    rpr = run._r.get_or_add_rPr()
    ea = rpr.find(qn("a:ea"))
    if ea is None:
        ea = etree.SubElement(rpr, qn("a:ea"))
    ea.set("typeface", FONT_NAME)


def add_textbox(
    slide: Slide,
    name: str,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    *,
    size: float = 16,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    color: RGBColor = BLACK,
    anchor: MSO_ANCHOR = MSO_ANCHOR.TOP,
) -> Slide:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.name = name
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.05)
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    set_run_font(run, size, bold=bold, color=color)
    return box


def build_cover_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 빈 레이아웃

    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(2.05), prs.slide_width, Inches(0.06)
    )
    bar.name = "cover_accent_bar"
    bar.fill.solid()
    bar.fill.fore_color.rgb = INDIGO
    bar.line.fill.background()

    add_textbox(
        slide,
        "cover_project",
        1.0,
        1.35,
        11.3,
        0.5,
        "사업명",
        size=18,
        align=PP_ALIGN.CENTER,
        color=INDIGO,
    )
    add_textbox(
        slide,
        "cover_title",
        1.0,
        2.4,
        11.3,
        1.3,
        "제안서",
        size=40,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    labels = [("발주처", "cover_client"), ("제안사", "cover_proposer"), ("제안일자", "cover_date")]
    for i, (label, name) in enumerate(labels):
        y = 4.7 + 0.5 * i
        add_textbox(
            slide,
            f"{name}_label",
            4.3,
            y,
            1.7,
            0.4,
            label,
            size=12,
            bold=True,
            align=PP_ALIGN.RIGHT,
            color=GRAY_TEXT,
        )
        add_textbox(slide, name, 6.2, y, 3.5, 0.4, "", size=12, align=PP_ALIGN.LEFT, color=SLATE)


def build_standard_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 상단 제목 바 (인디고 강조 막대 + 제목)
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.5), Inches(0.12), Inches(0.55)
    )
    accent.name = "title_accent"
    accent.fill.solid()
    accent.fill.fore_color.rgb = INDIGO
    accent.line.fill.background()

    add_textbox(
        slide,
        "slide_title",
        0.78,
        0.5,
        11.5,
        0.6,
        "슬라이드 제목",
        size=24,
        bold=True,
        anchor=MSO_ANCHOR.MIDDLE,
    )

    # 구분선
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.25), Inches(12.33), Pt(1)
    )
    line.name = "title_rule"
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
    line.line.fill.background()

    # 본문 불릿 영역 (렌더러가 불릿을 채움 — 서식 기준 문단 1개 포함)
    body = slide.shapes.add_textbox(Inches(0.78), Inches(1.55), Inches(11.8), Inches(5.4))
    body.name = "slide_body"
    tf = body.text_frame
    tf.word_wrap = True
    para = tf.paragraphs[0]
    run = para.add_run()
    run.text = "• 핵심 내용"
    set_run_font(run, 18, color=SLATE)


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    build_cover_slide(prs)
    build_standard_slide(prs)

    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT_PATH)
    print(f"생성 완료: {OUTPUT_PATH}")


if __name__ == "__main__":
    main()
