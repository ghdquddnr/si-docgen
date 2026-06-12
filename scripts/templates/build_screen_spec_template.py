"""화면정의서 PPT 양식(backend/templates/screen_spec.pptx) 제작 스크립트.

표지 슬라이드 + 화면정의 표준 슬라이드 1장을 구성하고, 렌더러가 식별할 수 있도록
모든 주입 대상 요소에 shape name 을 부여한다. 일회성 제작 도구이며 런타임 렌더러는
생성된 .pptx 파일만 사용한다.

실행: uv run python scripts/templates/build_screen_spec_template.py

shape name 규약:
- 표지: cover_title / cover_project / cover_system / cover_author / cover_date
- 표준 슬라이드: slide_title / screen_id / screen_name / menu_path
  / mockup_area / field_table / logic_text
"""

from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.shapes.base import BaseShape
from pptx.slide import Slide
from pptx.util import Inches, Pt

OUTPUT_PATH = Path(__file__).resolve().parents[2] / "backend" / "templates" / "screen_spec.pptx"

FONT_NAME = "맑은 고딕"
NAVY = RGBColor(0x1F, 0x3B, 0x57)
GRAY_FILL = RGBColor(0xD9, 0xD9, 0xD9)
LIGHT_FILL = RGBColor(0xF2, 0xF2, 0xF2)
BLACK = RGBColor(0x00, 0x00, 0x00)
GRAY_TEXT = RGBColor(0x80, 0x80, 0x80)

# "No Style, Table Grid" — 검정 격자선만 있는 기본 표 스타일
TABLE_GRID_STYLE_ID = "{2D5ABB26-0587-4C30-8999-92F81FD0307C}"

FIELD_TABLE_HEADERS = ["번호", "항목명", "유형", "필수", "설명"]
FIELD_TABLE_WIDTHS = [0.55, 1.4, 1.0, 0.55, 1.9]  # inch
SAMPLE_FIELD_ROW = ["①", "사용자 ID", "텍스트박스", "Y", "로그인 ID 입력"]


def set_run_font(run, size: float, *, bold: bool = False, color: RGBColor = BLACK) -> None:  # noqa: ANN001
    run.font.name = FONT_NAME
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    # 한글 폰트는 eastAsia(typeface)까지 지정해야 실제 적용된다
    rpr = run._r.get_or_add_rPr()
    ea = rpr.find(qn("a:ea"))
    if ea is None:
        ea = etree.SubElement(rpr, qn("a:ea"))
    ea.set("typeface", FONT_NAME)


def add_textbox(
    slide: Slide,
    name: str,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    *,
    size: float = 10,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    color: RGBColor = BLACK,
    fill: RGBColor | None = None,
    border: bool = False,
) -> BaseShape:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.name = name
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    set_run_font(run, size, bold=bold, color=color)
    if fill is not None:
        box.fill.solid()
        box.fill.fore_color.rgb = fill
    if border:
        box.line.color.rgb = BLACK
        box.line.width = Pt(0.75)
    return box


def set_table_grid_style(graphic_frame) -> None:  # noqa: ANN001
    """표 스타일을 'No Style, Table Grid'(검정 격자선)로 교체한다."""
    tbl = graphic_frame._element.graphic.graphicData.tbl
    tbl_pr = tbl.find(qn("a:tblPr"))
    style_id = tbl_pr.find(qn("a:tableStyleId"))
    if style_id is None:
        style_id = etree.SubElement(tbl_pr, qn("a:tableStyleId"))
    style_id.text = TABLE_GRID_STYLE_ID
    tbl_pr.set("firstRow", "0")
    tbl_pr.set("bandRow", "0")


def set_cell(
    cell,
    text: str,
    *,
    size: float = 10,
    bold: bool = False,
    center: bool = True,
    fill: RGBColor | None = None,
) -> None:  # noqa: ANN001
    cell.margin_left = cell.margin_right = Inches(0.04)
    cell.margin_top = cell.margin_bottom = Inches(0.02)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    para = cell.text_frame.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER if center else PP_ALIGN.LEFT
    run = para.add_run()
    run.text = text
    set_run_font(run, size, bold=bold)
    if fill is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill


def build_cover_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 빈 레이아웃

    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.18)
    )
    bar.name = "cover_accent_bar"
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()

    add_textbox(
        slide,
        "cover_title",
        1.67,
        2.3,
        10.0,
        1.1,
        "화면정의서",
        size=44,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide,
        "cover_project",
        1.67,
        3.6,
        10.0,
        0.6,
        "프로젝트명",
        size=20,
        align=PP_ALIGN.CENTER,
        color=NAVY,
    )

    labels = [("시스템명", "cover_system"), ("작성자", "cover_author"), ("작성일", "cover_date")]
    for i, (label, name) in enumerate(labels):
        y = 5.0 + 0.5 * i
        add_textbox(
            slide,
            f"{name}_label",
            4.5,
            y,
            1.5,
            0.4,
            label,
            size=12,
            bold=True,
            align=PP_ALIGN.RIGHT,
            color=GRAY_TEXT,
        )
        add_textbox(slide, name, 6.2, y, 3.2, 0.4, "", size=12, align=PP_ALIGN.LEFT)


def build_standard_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 상단 제목 + 문서명
    add_textbox(slide, "slide_title", 0.35, 0.15, 8.0, 0.5, "화면명", size=18, bold=True)
    add_textbox(
        slide,
        "doc_label",
        10.6,
        0.2,
        2.4,
        0.35,
        "화면정의서",
        size=10,
        align=PP_ALIGN.RIGHT,
        color=GRAY_TEXT,
    )

    # 화면 정보 스트립: 라벨(회색) + 값 텍스트 프레임
    strip = [
        ("화면 ID", "screen_id", 0.35, 1.0, 1.35, 1.7),
        ("화면명", "screen_name", 3.05, 1.0, 4.05, 2.6),
        ("메뉴 경로", "menu_path", 6.65, 1.0, 7.65, 5.3),
    ]
    for label, name, lx, lw, vx, vw in strip:
        add_textbox(
            slide,
            f"{name}_label",
            lx,
            0.8,
            lw,
            0.4,
            label,
            size=9,
            bold=True,
            align=PP_ALIGN.CENTER,
            fill=GRAY_FILL,
            border=True,
        )
        add_textbox(slide, name, vx, 0.8, vw, 0.4, "", size=9, border=True)

    # 목업 이미지 영역
    add_textbox(slide, "mockup_label", 0.35, 1.35, 3.0, 0.3, "■ 화면 목업", size=10, bold=True)
    mockup = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.35), Inches(1.7), Inches(7.0), Inches(5.4)
    )
    mockup.name = "mockup_area"
    mockup.fill.solid()
    mockup.fill.fore_color.rgb = LIGHT_FILL
    mockup.line.color.rgb = RGBColor(0xAA, 0xAA, 0xAA)
    mockup.line.width = Pt(0.75)
    tf = mockup.text_frame
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.add_run()
    run.text = "목업 이미지 영역"
    set_run_font(run, 12, color=GRAY_TEXT)

    # 항목 정의 표 (5열: 번호/항목명/유형/필수/설명, 헤더 + 서식 기준 행 1개)
    add_textbox(slide, "field_label", 7.55, 1.35, 3.0, 0.3, "■ 항목 정의", size=10, bold=True)
    total_w = sum(FIELD_TABLE_WIDTHS)
    frame = slide.shapes.add_table(
        2, len(FIELD_TABLE_HEADERS), Inches(7.55), Inches(1.7), Inches(total_w), Inches(0.62)
    )
    frame.name = "field_table"
    set_table_grid_style(frame)
    table = frame.table
    for col, width in zip(table.columns, FIELD_TABLE_WIDTHS, strict=True):
        col.width = Inches(width)
    for cell, header in zip(table.rows[0].cells, FIELD_TABLE_HEADERS, strict=True):
        set_cell(cell, header, bold=True, fill=GRAY_FILL, size=9)
    for i, (cell, sample) in enumerate(zip(table.rows[1].cells, SAMPLE_FIELD_ROW, strict=True)):
        set_cell(cell, sample, size=9, center=i != 4)
    table.rows[0].height = Inches(0.32)
    table.rows[1].height = Inches(0.3)

    # 처리 로직
    add_textbox(slide, "logic_label", 7.55, 4.85, 3.0, 0.3, "■ 처리 로직", size=10, bold=True)
    add_textbox(slide, "logic_text", 7.55, 5.2, total_w, 1.9, "", size=9, border=True)


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    build_cover_slide(prs)
    build_standard_slide(prs)

    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT_PATH)
    print(f"생성 완료: {OUTPUT_PATH}")


if __name__ == "__main__":
    main()
