"""제안서 LLM 생성 파이프라인 테스트 (LLM 모킹)."""

import json
from pathlib import Path
from typing import Any

import pytest
from pptx import Presentation

from app import cli
from app.pipelines.generate_proposal import generate_and_render_proposal, generate_proposal

ROOT = Path(__file__).resolve().parents[2]
INPUT = ROOT / "tests" / "fixtures" / "sources" / "sample_requirements.md"

MOCK_PROPOSAL: dict[str, Any] = {
    "project_name": "차세대 포털 사업",
    "system_name": "고객 포털",
    "author": "제안사",
    "client": "발주처",
    "written_date": "2026-06-16",
    "title": "차세대 포털 구축 제안서",
    "slides": [
        {"title": "사업 이해", "bullets": ["배경", "목표"]},
        {"title": "추진 전략", "bullets": ["전략1", "전략2"]},
        {"title": "기대 효과", "bullets": ["효과1"]},
    ],
}

COVER = {
    "project_name": "차세대 포털 사업",
    "system_name": "고객 포털",
    "author": "제안사",
    "client": "발주처",
    "written_date": "2026-06-16",
}


@pytest.fixture
def mock_llm(monkeypatch: pytest.MonkeyPatch) -> None:
    monkeypatch.setattr(
        "app.llm.generate.complete_json",
        lambda *a, **k: json.dumps(MOCK_PROPOSAL, ensure_ascii=False),
    )


def test_제안서_생성_및_검증(mock_llm: None) -> None:
    doc = generate_proposal(INPUT, **COVER)
    assert doc.title == "차세대 포털 구축 제안서"
    assert len(doc.slides) == 3
    assert doc.slides[0].title == "사업 이해"


def test_진행_콜백_단계_통지(mock_llm: None) -> None:
    stages: list[str] = []
    generate_proposal(INPUT, **COVER, on_progress=stages.append)
    assert stages == ["parsing", "generating"]


def test_generate_and_render(mock_llm: None, tmp_path: Path) -> None:
    result = generate_and_render_proposal(INPUT, tmp_path, **COVER)
    assert result.proposal_path.is_file()
    assert result.slide_count == 3
    assert result.bullet_count == 5  # 2 + 2 + 1
    # 표지 + 목차 + 내용 3 = 5 슬라이드
    prs = Presentation(str(result.proposal_path))
    assert len(prs.slides) == 5


def test_cli_proposal_종료코드_0(mock_llm: None, tmp_path: Path) -> None:
    code = cli.main(["proposal", "--input", str(INPUT), "--output", str(tmp_path)])
    assert code == 0
    assert (tmp_path / "proposal.pptx").is_file()
