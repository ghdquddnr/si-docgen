"""제안서 PPT 렌더러 골든 테스트 — 고정 JSON → pptx 텍스트 비교."""

import json
from pathlib import Path

import pytest
from pptx import Presentation

from app.renderers.proposal_renderer import render_proposal
from app.schemas.proposal import ProposalDocument

ROOT = Path(__file__).resolve().parents[2]
TEMPLATE = ROOT / "backend" / "templates" / "proposal.pptx"
FIXTURE = ROOT / "tests" / "golden" / "fixtures" / "proposal_8.json"


@pytest.fixture
def doc() -> ProposalDocument:
    return ProposalDocument.model_validate(json.loads(FIXTURE.read_text(encoding="utf-8")))


def _shape_text(slide, name: str) -> str:
    for shape in slide.shapes:
        if shape.name == name and shape.has_text_frame:
            return shape.text_frame.text
    raise AssertionError(f"shape '{name}' 없음")


def test_슬라이드_수_표지_목차_내용(doc: ProposalDocument, tmp_path: Path) -> None:
    out = render_proposal(doc, TEMPLATE, tmp_path / "p.pptx")
    prs = Presentation(str(out))
    # 표지 1 + 목차 1 + 내용 7
    assert len(prs.slides) == 2 + len(doc.slides) == 9


def test_목차_슬라이드_자동생성(doc: ProposalDocument, tmp_path: Path) -> None:
    out = render_proposal(doc, TEMPLATE, tmp_path / "p.pptx")
    toc = list(Presentation(str(out)).slides)[1]
    assert _shape_text(toc, "slide_title") == "목차"
    body = _shape_text(toc, "slide_body")
    # 모든 섹션 제목이 번호와 함께 목차에 나열됨
    for i, s in enumerate(doc.slides, start=1):
        assert f"{i}. {s.title}" in body


def test_표지_정보_주입(doc: ProposalDocument, tmp_path: Path) -> None:
    out = render_proposal(doc, TEMPLATE, tmp_path / "p.pptx")
    cover = Presentation(str(out)).slides[0]
    assert _shape_text(cover, "cover_title") == "차세대 고객 포털 구축 제안서"
    assert _shape_text(cover, "cover_client") == "한국전력공사"
    assert _shape_text(cover, "cover_proposer") == "(주)에스아이파트너스"
    assert _shape_text(cover, "cover_date") == "2026-06-16"


def test_내용_슬라이드_제목과_불릿(doc: ProposalDocument, tmp_path: Path) -> None:
    out = render_proposal(doc, TEMPLATE, tmp_path / "p.pptx")
    # 표지(0)·목차(1) 다음부터 내용 슬라이드
    content = list(Presentation(str(out)).slides)[2:]
    assert [_shape_text(s, "slide_title") for s in content] == [s.title for s in doc.slides]

    # 첫 내용 슬라이드의 본문에 모든 불릿이 포함됨
    body = _shape_text(content[0], "slide_body")
    for bullet in doc.slides[0].bullets:
        assert bullet in body
    assert "•" in body  # 불릿 머리표


def test_표준_슬라이드_원본은_제거됨(doc: ProposalDocument, tmp_path: Path) -> None:
    # 템플릿의 플레이스홀더 '슬라이드 제목' 이 결과물에 남지 않는다
    out = render_proposal(doc, TEMPLATE, tmp_path / "p.pptx")
    titles = [_shape_text(s, "slide_title") for s in list(Presentation(str(out)).slides)[1:]]
    assert "슬라이드 제목" not in titles


def test_libreoffice_없이_pptx_열림(doc: ProposalDocument, tmp_path: Path) -> None:
    out = render_proposal(doc, TEMPLATE, tmp_path / "p.pptx")
    prs = Presentation(str(out))  # 손상 없이 열리는지
    assert prs.slide_width is not None
