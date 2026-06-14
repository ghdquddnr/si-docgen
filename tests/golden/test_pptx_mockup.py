"""화면정의서 PPT 목업 — 편집 가능한 도형 렌더링 골든 테스트.

목업이 이미지가 아니라 항목별 PPT 도형(번호 배지·라벨·컨트롤)으로 그려지는지 확인한다.
"""

import json
from pathlib import Path
from typing import Any

import pytest
from pptx import Presentation
from pptx.shapes.base import BaseShape

from app.renderers.pptx_renderer import render_screen_spec
from app.schemas.screen_spec import ScreenSpecDocument

ROOT = Path(__file__).resolve().parents[2]
TEMPLATE = ROOT / "backend" / "templates" / "screen_spec.pptx"
FIXTURE = ROOT / "tests" / "golden" / "fixtures" / "screen_spec_5.json"


@pytest.fixture(scope="module")
def fixture_data() -> dict[str, Any]:
    return json.loads(FIXTURE.read_text(encoding="utf-8"))


@pytest.fixture(scope="module")
def rendered(fixture_data: dict[str, Any], tmp_path_factory: pytest.TempPathFactory) -> Path:
    doc = ScreenSpecDocument.model_validate(fixture_data)
    out = tmp_path_factory.mktemp("pptx") / "screen_spec.pptx"
    return render_screen_spec(doc, TEMPLATE, out)


def _named(slide, prefix: str) -> dict[str, BaseShape]:
    return {s.name: s for s in slide.shapes if s.name.startswith(prefix)}


def test_화면당_슬라이드_생성(rendered: Path, fixture_data: dict[str, Any]) -> None:
    prs = Presentation(str(rendered))
    # 표지 1 + 화면 수
    assert len(prs.slides) == 1 + len(fixture_data["screens"])


def test_목업이_도형으로_그려짐(rendered: Path, fixture_data: dict[str, Any]) -> None:
    prs = Presentation(str(rendered))
    first_screen = fixture_data["screens"][0]
    slide = prs.slides[1]  # 표지(0) 다음이 첫 화면

    badges = _named(slide, "mockup_badge_")
    labels = _named(slide, "mockup_label_")
    ctrls = _named(slide, "mockup_ctrl_")

    n = len(first_screen["fields"])
    assert len(badges) == n
    assert len(labels) == n
    assert len(ctrls) == n

    # 이미지(그림) 도형은 없어야 한다 (편집 가능 도형만)
    assert not any(s.shape_type == 13 for s in slide.shapes)  # 13 = PICTURE


def test_배지_라벨_텍스트_일치(rendered: Path, fixture_data: dict[str, Any]) -> None:
    prs = Presentation(str(rendered))
    first_screen = fixture_data["screens"][0]
    slide = prs.slides[1]

    for field in first_screen["fields"]:
        badge = next(s for s in slide.shapes if s.name == f"mockup_badge_{field['no']}")
        label = next(s for s in slide.shapes if s.name == f"mockup_label_{field['no']}")
        assert badge.text_frame.text == str(field["no"])
        assert field["name"] in label.text_frame.text


def test_항목_없는_화면은_컨트롤_없음(tmp_path: Path) -> None:
    doc = ScreenSpecDocument.model_validate(
        {
            "project_name": "P",
            "system_name": "S",
            "author": "A",
            "written_date": "2026-06-14",
            "screens": [
                {
                    "screen_id": "SCR-001",
                    "screen_name": "빈 화면",
                    "menu_path": "홈",
                    "fields": [
                        {"no": 1, "name": "항목", "field_type": "텍스트박스", "required": False}
                    ],
                }
            ],
        }
    )
    out = render_screen_spec(doc, TEMPLATE, tmp_path / "out.pptx")
    slide = Presentation(str(out)).slides[1]
    assert len([s for s in slide.shapes if s.name.startswith("mockup_ctrl_")]) == 1
